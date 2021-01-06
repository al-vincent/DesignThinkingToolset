from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from azure.storage.blob import BlobServiceClient
from azure.core.exceptions import ResourceExistsError, ResourceNotFoundError

import os
from datetime import datetime, timedelta
from PIL import Image
from io import BytesIO
import logging

# Get a logger instance
logger = logging.getLogger(__name__)

# ================================================================================================
# CREATE POWERPOINT
# ================================================================================================
class SnipPptxCreator:
    def __init__(self, image_filename, image_bytes, text_boxes):
        self.image_filename = image_filename
        self.image_bytes = image_bytes
        self.image = self.get_image()
        self.text_boxes = text_boxes
        self.pres = Presentation()

    def get_image(self):
        """
        Open the image using pillow for easy manipulation
        """
        try:
            return Image.open(self.image_bytes)
        except Exception as err:
            logger.error(f"An exception occurred. Sys error: {err}")
            # --------------------------------
            # TODO: change the return! If this happens, it's pretty critical...!
            # --------------------------------
            return None

    def add_title_slide(self):
        # set the layout and create the slide
        title_slide_layout = self.pres.slide_layouts[0]
        slide = self.pres.slides.add_slide(title_slide_layout)

        # create the title and subtitle objects
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        # add the title and subtitle
        title.text = self.image_filename
        now = datetime.now()
        subtitle.text = f"Text Analysis Output \nCreated: {now.strftime('%d %b %Y, %H:%M')}"

    def add_original_image_slide(self):        
        """
        Add a slide with a title and the original image embedded
        """
        # add a blank slide to the presentation
        title_only_slide_layout = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(title_only_slide_layout)

        # create the title and add its text
        title = slide.shapes.title
        title.text = "Original Image"      

        # landscape image
        if self.image.width > self.image.height:
            # set the width of the image to be 70% of the slide-width
            percent = 0.8
            width = int(self.pres.slide_width * percent)
            # calculate the new image height
            height = int((width / self.image.width) * self.image.height)
            # ensure that the bottom of the image is above the bottom of the slide
            while (height + title.top + title.height > self.pres.slide_height):
                percent -= 0.05
                width = int(self.pres.slide_width * percent)
                height = int((width / self.image.width) * self.image.height)
            
            # centre the image horizontally, and place just below the title
            left = int((self.pres.slide_width - width) / 2)
            top = title.top + title.height
            # record the image size and position; we'll need it later!
            self.image_size = {"left": left, "top": top, "width": width, "height": height}
            # add the image to the slide
            slide.shapes.add_picture(self.image_bytes, left, top, width=width)
        else:
            # set the height of the image to be 80% of <space-below-bottom-of-title>
            height = int((self.pres.slide_height - (title.top + title.height)) * 0.8)
            # calculate the new image width
            width = int((height / self.image.height) * self.image.width)

            # centre the image horizontally, and place just below the title
            left = int((self.pres.slide_width - width) / 2)
            top = title.top + title.height
            # record the image size and position; we'll need it later!
            self.image_size = {"left": left, "top": top, "width": width, "height": height}            
            # add the image to the slide
            slide.shapes.add_picture(self.image_bytes, left, top, width=width, height=height)

    def add_image_and_text_slide(self):
        """
        Add a slide with a title and the original image embedded, and text boxes
        overlaid onto the image
        """
        # add a blank slide to the presentation
        title_only_slide_layout = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(title_only_slide_layout)

        # create the title and add its text
        title = slide.shapes.title
        title.text = "Original Image with Extracted Text"
        
        # add the image to the slide
        slide.shapes.add_picture(self.image_bytes, self.image_size["left"], self.image_size["top"], 
                                self.image_size["width"], self.image_size["height"])

        # overlay the textboxes on the image
        for text_box in self.text_boxes:            
            self.add_textbox(slide, text_box)
    
    # ***************************************************************************
    # NOTE: python-pptx doesn't implement shape fill transparency yet (Aug 2020). 
    # The only way to manipulate it (and it will help a lot here) is to 
    # manipulate the underlying XML of the PowerPoint, which is...tricky.
    # 
    # Thankfully, stackoverflow has the answer in the two methods below
    # https://stackoverflow.com/a/57258217
    # ***************************************************************************
    def sub_element(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def set_shape_transparency(self, shape, alpha):
        """ 
        Set the transparency (alpha) of a shape
        """
        ts = shape.fill._xPr.solidFill
        sf = ts.get_or_change_to_srgbClr()
        self.sub_element(sf, 'a:alpha', val=str(alpha))

    def add_text_as_bullets_slide(self):
        """
        Add a slide with a title and the text from each region 
        as a bulleted list
        """
        # add a blank slide to the presentation
        title_and_content_slide_layout = self.pres.slide_layouts[1]
        slide = self.pres.slides.add_slide(title_and_content_slide_layout)

        # create the title and add its text
        title = slide.shapes.title
        title.text = "Text Extracted From Each Region"

        # add text to the body
        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.text = "Each bullet shows all text in a single region:"
        for text_box in self.text_boxes:
            para = text_frame.add_paragraph()
            para.text = text_box["text"]
            para.level = 1

    def add_textbox(self, slide, text_box):
        # convert the coords from normalised in range [0,1] to image coords
        left = self.image_size["left"] + (self.image_size["width"] * text_box["x"])
        top = self.image_size["top"] + (self.image_size["height"] * text_box["y"])
        width = self.image_size["width"] * text_box["width"]
        height = self.image_size["height"] * text_box["height"]

        # add the text-box
        txt_box = slide.shapes.add_textbox(left, top, width, height)

        # make the fill solid white, and change the transparency
        fill = txt_box.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        # 70,000 represents 70% *OPACITY* (i.e. 30% transparency)
        self.set_shape_transparency(txt_box,70000)

        # add the text and format it
        text_frame = txt_box.text_frame
        text_frame.text = text_box["text"]
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE        
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        # ------------------------------------------------------------------
        # NOTE: the lines below allow the font-size to be changed; *BUT* the 
        # TEXT_TO_FIT_SHAPE sizing doesn't then work. So ignoring.
        # ------------------------------------------------------------------
        # run = para.add_run()
        # run.text = text_box["text"]
        # font = run.font
        # font.size = Pt(80)
        # ------------------------------------------------------------------

    def save_pres_as_bytes(self):
        out = BytesIO()
        self.pres.save(out)
        return out.getvalue()

    def build_pres(self):
        self.add_title_slide()
        self.add_original_image_slide()
        self.add_image_and_text_slide()
        self.add_text_as_bullets_slide()
        return self.save_pres_as_bytes()

# ================================================================================================
# SAVE FILE TO BLOB STORAGE
# ================================================================================================
class SaveFileToAzureBlobStorage:
    def __init__(self, connect_str, container_name):
        self.connect_str = connect_str
        self.container_name = container_name

    def create_blob_service_client(self):
        # Create the BlobServiceClient object which will be used to create a container client
        try:
            logger.info("The BlobServiceClient was created successfully")
            return BlobServiceClient.from_connection_string(self.connect_str)
        except ValueError as err:
            # Occurs when the connection string is either malformed or empty
            logger.error(f"The connection string '{self.connect_str}' is invalid; the BlobServiceClient was not created. Sys error: {err}")
            return None
        except AttributeError as err:
            # Occurs when the connection string is None (for e.g.)
            logger.error(f"The connection string '{self.connect_str}' is invalid; the BlobServiceClient was not created. Sys error: {err}")
            return None
        except Exception as err:
            logger.error(f"Another error occurred; the BlobServiceClient was not created. Sys error: {err}")
            return None

    def create_container(self):
        blob_service_client = self.create_blob_service_client()
        
        if blob_service_client is not None:
            # Create the container, with blobs (*NOT* the container) being publicly accessible
            try:
                logger.info(f"The container {self.container_name} was created successfully")
                return blob_service_client.create_container(self.container_name, public_access="blob")        
            except ResourceExistsError as err:
                # this is raised if the container already exists
                logger.error(f"A container named {self.container_name} already exists. Sys error: {err}")
                return None
            except Exception as err:
                logger.error(f"Another error occurred; failed to create the container. Sys error: {err}")
                return None
        else:
            return None

    def create_blob_client(self, file_name):
        blob_service_client = self.create_blob_service_client()
        
        if blob_service_client is not None:            
            try:
                logger.info(f"The blob service client was created successfully")
                return blob_service_client.get_blob_client(container=self.container_name,
                                                        blob=file_name)
            except Exception as err:
                logger.error(f"Failed to create the blob. Sys error: {err}")
                return None
        else:
            return None

    def upload_bytes_to_container(self, file_name, file_bytes):
        blob_client = self.create_blob_client(file_name=file_name)

        if blob_client is not None:
            try:
                blob_client.upload_blob(file_bytes)
            except Exception as err:
                logger.error(f"Another exception occurred; file {file_name} was NOT uploaded successfully. Sys error: {err}")
                return None
            else:
                logger.info(f"File {file_name} uploaded successfully")
                return blob_client.url
        else: 
            logger.error(f"Blob client creation failed; blob {file_name} was NOT uploaded successfully")
            return None
        
    def create_container_and_upload_file(self, file_name, file_bytes):
        container = self.create_container()
        
        return self.upload_bytes_to_container(file_name, file_bytes) if container else None

    def list_blobs_in_container(self):
        # NOTE: this is primarily useful as a check that uploads have completed 
        # successfully (or otherwise)
        pass

    def delete_container(self):
        pass

    def delete_blob(self, file_name):
        pass


# ================================================================================================
# DRIVER
# ================================================================================================
def main():
    infile = "match_words_regions_1.jpg"
    current_dir = os.path.dirname(os.path.abspath(__file__))
    image_path = os.path.join(current_dir, "tests", "resources", "test_images", infile)
    outfile = "test.pptx"
    image_bytes = open(image_path, "rb")

    # -----------------------------------------------------------------------------
    # create the presentation, save in /media/tmp
    # -----------------------------------------------------------------------------
    text_boxes = [
        {"x": 0.00842, "y": 0.03377, "width": 0.23617, "height": 0.34575, "text": "Stay Home"},
        {"x": 0.48438, "y": 0.04555, "width": 0.24382, "height": 0.33621, "text": "KEEP YOUR DISTANCE"},
        {"x": 0.25384, "y": 0.04926, "width": 0.22724, "height": 0.35479, "text": "Be Safe"},
        {"x": 0.74035, "y": 0.05551, "width": 0.23943, "height": 0.34818, "text": "Don't out !"},
        {"x": 0.01450, "y": 0.43439, "width": 0.23028, "height": 0.36581, "text": "WRITE A LETTER"},
        {"x": 0.24843, "y": 0.43413, "width": 0.23775, "height": 0.36817, "text": "Read a book"},
        {"x": 0.48965, "y": 0.43800, "width": 0.23516, "height": 0.36250, "text": "Online Chat"},
        {"x": 0.74491, "y": 0.44507, "width": 0.24687, "height": 0.35204, "text": "Phone a friend"}
    ]

    spc = SnipPptxCreator(image_filename=infile,
                        image_bytes=image_bytes,
                        text_boxes=text_boxes)
    pres_bytes = spc.build_pres()
    print(f"Presentation created")

    # -----------------------------------------------------------------------------
    # store the presentation in blob storage
    # -----------------------------------------------------------------------------
    connect_str = os.getenv('SNIP_BLOB_STORAGE_CONN_STR')
    container_name = "snip-test-container-01"
    
    sftabs = SaveFileToAzureBlobStorage(connect_str=connect_str, 
                                        container_name=container_name)
    container_created = sftabs.create_container()
    if container_created:
        url = sftabs.upload_bytes_to_container(file_name=outfile, file_bytes=pres_bytes)
        if url is not None:
            print(f"File saved to blob storage; see {url}")
        else:
            print(f"The file was NOT saved successfully")
    else:
        print(f"The container was NOT created successfully")

    # url = sftabs.upload_bytes_to_blob_storage(infile, image_bytes)
    # if url is not None:
    #     print(f"File saved to blob storage; see {url}")
    # else:
    #     print(f"The file was NOT saved successfully")
    
if __name__ == "__main__":
    main()